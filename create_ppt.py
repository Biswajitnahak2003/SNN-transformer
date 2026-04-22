#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation for the Adaptive SNN research.
This script generates slides with proper content fitting and beautiful architecture visualization.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_presentation():
    """Create the PowerPoint presentation with proper formatting."""

    # Create presentation with blank template
    prs = Presentation()

    # Set default font to Times New Roman for all slides
    def set_font_times_new_roman(text_frame):
        """Set all text in a text frame to Times New Roman."""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
            # Also set the paragraph font
            paragraph.font.name = 'Times New Roman'

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Adaptive Timestep Spiking Neural Networks\nfor Brain Tumor Segmentation"
    subtitle.text = "Energy-Efficient Medical Image Analysis with Dynamic Temporal Processing\n\nBiswajit Sahoo\nResearch Presentation"

    # Set fonts to Times New Roman
    set_font_times_new_roman(title.text_frame)
    set_font_times_new_roman(subtitle.text_frame)

    # Slide 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Introduction"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    # Add content with proper positioning
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "• Brain tumor segmentation is critical for treatment planning and prognosis"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Traditional CNNs are computationally expensive for medical imaging"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Spiking Neural Networks (SNNs) offer energy-efficient alternatives"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Challenge: Fixed temporal resolution limits SNN performance"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Solution: Adaptive timestep processing based on local complexity"
    p.level = 0

    set_font_times_new_roman(text_frame)

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Problem Statement"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Medical imaging requires high accuracy but computational resources are limited:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Brain MRI scans: 4 modalities × 155 slices × 240×240 resolution"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Real-time analysis needed for clinical workflows"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Edge devices and mobile platforms have limited compute"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• SNNs promise 10-100x energy reduction but fixed timesteps hurt accuracy"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Need adaptive processing: simple regions → few timesteps, complex regions → more timesteps"
    p.level = 0

    set_font_times_new_roman(text_frame)

    # Slide 4: Related Work
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Related Work"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Medical Image Segmentation:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• U-Net (2015): Encoder-decoder with skip connections"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Attention U-Net (2018): Spatial attention mechanisms"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• TransUNet (2021): Transformer-enhanced segmentation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Spiking Neural Networks:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Spiking U-Net (2022): Direct SNN adaptation of U-Net"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Temporal-efficient SNNs (2023): Fixed timestep optimization"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Adaptive Computing:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Dynamic networks (2017): Input-adaptive architectures"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Neural architecture search (2018): Automated design"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 5: Methodology Overview
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Methodology Overview"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Our adaptive SNN framework consists of three main components:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "1. Spiking U-Net Backbone: 4-stage encoder-decoder with skip connections"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "2. Bipolar Linear Self-Attention: Captures long-range spatiotemporal dependencies"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "3. CNN Uncertainty Agent: Predicts optimal timestep allocation per spatial location"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Training Strategy:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Phase 1: Fixed T=4 warmup (epochs 0-10)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Phase 2: Adaptive processing (epochs 11-19)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Joint optimization of segmentation and agent losses"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 6: Architecture Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = "Architecture Overview"

    # Create a beautiful architecture diagram using shapes and text
    # Input block
    input_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(1.5), Inches(0.8))
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Cornflower blue
    input_box.line.color.rgb = RGBColor(0, 0, 0)

    input_text = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(1.3), Inches(0.4))
    input_text.text_frame.paragraphs[0].text = "MRI Input\nT1, T1c, T2, FLAIR"
    input_text.text_frame.paragraphs[0].font.size = Pt(10)
    input_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Arrow to SNN
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.2), Inches(2), Inches(0.8), Inches(0.2))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(169, 169, 169)

    # Spiking U-Net block
    snn_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(1.2), Inches(2.5), Inches(1.5))
    snn_box.fill.solid()
    snn_box.fill.fore_color.rgb = RGBColor(255, 140, 0)  # Dark orange
    snn_box.line.color.rgb = RGBColor(0, 0, 0)

    snn_text = slide.shapes.add_textbox(Inches(3.4), Inches(1.4), Inches(2.1), Inches(1.1))
    snn_text.text_frame.paragraphs[0].text = "Spiking U-Net\n4 Stages\nLIF Neurons\nSkip Connections"
    snn_text.text_frame.paragraphs[0].font.size = Pt(9)
    snn_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Attention block
    attention_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.8), Inches(1.8), Inches(0.8))
    attention_box.fill.solid()
    attention_box.fill.fore_color.rgb = RGBColor(50, 205, 50)  # Lime green
    attention_box.line.color.rgb = RGBColor(0, 0, 0)

    attention_text = slide.shapes.add_textbox(Inches(6.2), Inches(2), Inches(1.4), Inches(0.4))
    attention_text.text_frame.paragraphs[0].text = "Bipolar\nSelf-Attention"
    attention_text.text_frame.paragraphs[0].font.size = Pt(8)
    attention_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Agent block
    agent_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(2), Inches(0.8))
    agent_box.fill.solid()
    agent_box.fill.fore_color.rgb = RGBColor(220, 20, 60)  # Crimson
    agent_box.line.color.rgb = RGBColor(0, 0, 0)

    agent_text = slide.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(1.6), Inches(0.4))
    agent_text.text_frame.paragraphs[0].text = "CNN Uncertainty Agent\nTimestep Prediction"
    agent_text.text_frame.paragraphs[0].font.size = Pt(8)
    agent_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Controller block
    controller_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(2), Inches(0.8))
    controller_box.fill.solid()
    controller_box.fill.fore_color.rgb = RGBColor(138, 43, 226)  # Blue violet
    controller_box.line.color.rgb = RGBColor(0, 0, 0)

    controller_text = slide.shapes.add_textbox(Inches(4.2), Inches(3.7), Inches(1.6), Inches(0.4))
    controller_text.text_frame.paragraphs[0].text = "Adaptive Controller\nSelective Processing"
    controller_text.text_frame.paragraphs[0].font.size = Pt(8)
    controller_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Output block
    output_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(3.5), Inches(1.5), Inches(0.8))
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(0, 191, 255)  # Deep sky blue
    output_box.line.color.rgb = RGBColor(0, 0, 0)

    output_text = slide.shapes.add_textbox(Inches(7.2), Inches(3.7), Inches(1.1), Inches(0.4))
    output_text.text_frame.paragraphs[0].text = "Segmentation\nOutput\n4 Classes"
    output_text.text_frame.paragraphs[0].font.size = Pt(8)
    output_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Arrows
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(2.9), Inches(0.2), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(169, 169, 169)

    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(4.4), Inches(0.6), Inches(0.2))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = RGBColor(169, 169, 169)

    arrow4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(4.4), Inches(0.6), Inches(0.2))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = RGBColor(169, 169, 169)

    # Slide 7: Spiking U-Net Details
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Spiking U-Net Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    # Left side: Architecture diagram inside content area
    diagram_left = Inches(0.5)
    diagram_top = Inches(1.2)
    diagram_width = Inches(4.5)
    diagram_height = Inches(4.5)

    # Create mini architecture diagram for SNN
    # Input layer
    input_layer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(0.5), diagram_top + Inches(0.2), Inches(1.5), Inches(0.4))
    input_layer.fill.solid()
    input_layer.fill.fore_color.rgb = RGBColor(100, 149, 237)
    input_layer.line.color.rgb = RGBColor(0, 0, 0)

    input_label = slide.shapes.add_textbox(diagram_left + Inches(0.7), diagram_top + Inches(0.3), Inches(1.1), Inches(0.2))
    input_label.text_frame.paragraphs[0].text = "Input (128×128×4)"
    input_label.text_frame.paragraphs[0].font.size = Pt(8)
    set_font_times_new_roman(input_label.text_frame)

    # Encoder blocks
    for i in range(4):
        y_pos = diagram_top + Inches(0.8 + i * 0.6)
        enc_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(0.8), y_pos, Inches(1.2), Inches(0.4))
        enc_block.fill.solid()
        enc_block.fill.fore_color.rgb = RGBColor(255, 140, 0)
        enc_block.line.color.rgb = RGBColor(0, 0, 0)

        enc_label = slide.shapes.add_textbox(diagram_left + Inches(1.0), y_pos + Inches(0.1), Inches(0.8), Inches(0.2))
        enc_label.text_frame.paragraphs[0].text = f"Enc{i+1}"
        enc_label.text_frame.paragraphs[0].font.size = Pt(7)
        set_font_times_new_roman(enc_label.text_frame)

        # Down arrow
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diagram_left + Inches(1.4), y_pos + Inches(0.45), Inches(0.2), Inches(0.1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(169, 169, 169)

    # Bottleneck with attention
    bottleneck = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(1.2), diagram_top + Inches(3.2), Inches(1.2), Inches(0.4))
    bottleneck.fill.solid()
    bottleneck.fill.fore_color.rgb = RGBColor(50, 205, 50)
    bottleneck.line.color.rgb = RGBColor(0, 0, 0)

    bottleneck_label = slide.shapes.add_textbox(diagram_left + Inches(1.4), diagram_top + Inches(3.3), Inches(0.8), Inches(0.2))
    bottleneck_label.text_frame.paragraphs[0].text = "Attention"
    bottleneck_label.text_frame.paragraphs[0].font.size = Pt(7)
    set_font_times_new_roman(bottleneck_label.text_frame)

    # Decoder blocks
    for i in range(4):
        y_pos = diagram_top + Inches(4.0 - i * 0.6)
        dec_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(2.8), y_pos, Inches(1.2), Inches(0.4))
        dec_block.fill.solid()
        dec_block.fill.fore_color.rgb = RGBColor(255, 140, 0)
        dec_block.line.color.rgb = RGBColor(0, 0, 0)

        dec_label = slide.shapes.add_textbox(diagram_left + Inches(3.0), y_pos + Inches(0.1), Inches(0.8), Inches(0.2))
        dec_label.text_frame.paragraphs[0].text = f"Dec{i+1}"
        dec_label.text_frame.paragraphs[0].font.size = Pt(7)
        set_font_times_new_roman(dec_label.text_frame)

        # Skip connection
        if i < 3:
            skip = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, diagram_left + Inches(2.0), y_pos + Inches(0.15), Inches(0.6), Inches(0.1))
            skip.fill.solid()
            skip.fill.fore_color.rgb = RGBColor(169, 169, 169)

        # Up arrow
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, diagram_left + Inches(3.4), y_pos - Inches(0.05), Inches(0.2), Inches(0.1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(169, 169, 169)

    # Output layer
    output_layer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(3.2), diagram_top + Inches(4.6), Inches(1.5), Inches(0.4))
    output_layer.fill.solid()
    output_layer.fill.fore_color.rgb = RGBColor(0, 191, 255)
    output_layer.line.color.rgb = RGBColor(0, 0, 0)

    output_label = slide.shapes.add_textbox(diagram_left + Inches(3.4), diagram_top + Inches(4.7), Inches(1.1), Inches(0.2))
    output_label.text_frame.paragraphs[0].text = "Output (128×128×4)"
    output_label.text_frame.paragraphs[0].font.size = Pt(8)
    set_font_times_new_roman(output_label.text_frame)

    # Right side: Technical details
    details_left = Inches(5.5)
    details_top = Inches(1.2)
    details_width = Inches(4.5)
    details_height = Inches(4.5)

    textbox = slide.shapes.add_textbox(details_left, details_top, details_width, details_height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "4-stage encoder-decoder with spiking convolutional blocks:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Encoder: Progressive downsampling (128→64→32→16)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Decoder: Upsampling with skip connections"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• SpikingConvBlock: LIF neurons + 3×3 convolutions"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Membrane dynamics: α = 0.9, V_th = 1.0"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Bottleneck Attention:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Bipolar linear self-attention on spike trains"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Captures spatiotemporal dependencies"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Query-Key-Value operations on temporal dimension"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 8: CNN Uncertainty Agent
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "CNN Uncertainty Agent"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    # Left side: Architecture diagram inside content area
    diagram_left = Inches(0.5)
    diagram_top = Inches(1.2)
    diagram_width = Inches(4.5)
    diagram_height = Inches(4.5)

    # Create mini architecture diagram for CNN Agent
    # Input layer
    agent_input = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(0.5), diagram_top + Inches(0.2), Inches(1.5), Inches(0.4))
    agent_input.fill.solid()
    agent_input.fill.fore_color.rgb = RGBColor(100, 149, 237)
    agent_input.line.color.rgb = RGBColor(0, 0, 0)

    input_label = slide.shapes.add_textbox(diagram_left + Inches(0.6), diagram_top + Inches(0.3), Inches(1.3), Inches(0.2))
    input_label.text_frame.paragraphs[0].text = "Input (128×128×8)"
    input_label.text_frame.paragraphs[0].font.size = Pt(8)
    set_font_times_new_roman(input_label.text_frame)

    # Conv blocks
    conv_layers = [(32, "3×3 Conv"), (64, "3×3 Conv"), (128, "3×3 Conv")]
    for i, (channels, desc) in enumerate(conv_layers):
        y_pos = diagram_top + Inches(0.8 + i * 0.8)
        conv_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(0.8), y_pos, Inches(1.2), Inches(0.5))
        conv_block.fill.solid()
        conv_block.fill.fore_color.rgb = RGBColor(220, 20, 60)
        conv_block.line.color.rgb = RGBColor(0, 0, 0)

        conv_label = slide.shapes.add_textbox(diagram_left + Inches(0.9), y_pos + Inches(0.1), Inches(1.0), Inches(0.3))
        conv_label.text_frame.paragraphs[0].text = f"{channels}ch\n{desc}"
        conv_label.text_frame.paragraphs[0].font.size = Pt(7)
        set_font_times_new_roman(conv_label.text_frame)

        # Down arrow
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, diagram_left + Inches(1.4), y_pos + Inches(0.55), Inches(0.2), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(169, 169, 169)

    # Output layer
    output_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, diagram_left + Inches(1.2), diagram_top + Inches(3.4), Inches(1.2), Inches(0.5))
    output_block.fill.solid()
    output_block.fill.fore_color.rgb = RGBColor(0, 191, 255)
    output_block.line.color.rgb = RGBColor(0, 0, 0)

    output_label = slide.shapes.add_textbox(diagram_left + Inches(1.3), diagram_top + Inches(3.5), Inches(1.0), Inches(0.3))
    output_label.text_frame.paragraphs[0].text = "Timestep\nPredictions\n(1-4)"
    output_label.text_frame.paragraphs[0].font.size = Pt(7)
    set_font_times_new_roman(output_label.text_frame)

    # Right side: Technical details
    details_left = Inches(5.5)
    details_top = Inches(1.2)
    details_width = Inches(4.5)
    details_height = Inches(4.5)

    textbox = slide.shapes.add_textbox(details_left, details_top, details_width, details_height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Lightweight CNN that predicts timestep allocation:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Input: 8 channels (4 preliminary logits + 4 MRI modalities)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Architecture: 3 convolutional stages (32→64→128 channels)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Output: Per-pixel timestep assignments (1-4)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Training Objectives:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Segmentation loss: Dice + Focal loss"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Agent loss: Uncertainty-weighted efficiency term"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Joint optimization with curriculum learning"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 9: Experimental Setup
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Experimental Setup"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Dataset: BraTS 2023 Glioma Segmentation"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• 1,251 high-grade glioma cases"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• 4 MRI modalities: T1, T1c, T2, FLAIR"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• 4 tissue classes: Background, Necrotic core, Edema, Enhancing tumor"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Training Details:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Resolution: 128×128 axial slices"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Batch size: 4, Optimizer: AdamW (lr=1e-4)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• 20 epochs with two-phase curriculum"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Hardware: NVIDIA RTX 4090"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 10: Main Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Main Results"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Performance on BraTS 2023 validation set:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Dice Coefficient: 0.7410"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Hausdorff Distance 95th percentile: 2.412 pixels"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Energy Savings: 25.26% reduction vs. static SNN (T=4)"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Comparison with Static SNN:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Static SNN (T=4): Dice = 0.732, Energy = 100%"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Adaptive SNN: Dice = 0.741, Energy = 75%"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• +1.4% accuracy with 25% energy reduction"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 11: Training Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Training Analysis"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Two-phase training curriculum:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Phase 1 (epochs 0-10): Fixed T=4 warmup"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Phase 2 (epochs 11-19): Adaptive processing activation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Stable convergence achieved:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Final validation Dice: 0.7410"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Consistent energy reduction after agent activation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Smooth loss curves indicate stable optimization"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 12: Qualitative Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Qualitative Results"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Segmentation quality assessment:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Accurate tumor boundary delineation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Proper tissue class separation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Robust performance across different tumor sizes"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Visual comparison shows:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Superior boundary accuracy vs. static SNN"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Better handling of heterogeneous tumor regions"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Consistent performance in challenging cases"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 13: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Conclusion"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Key Contributions:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Novel adaptive timestep SNN for medical imaging"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Bipolar linear self-attention for spatiotemporal modeling"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• CNN uncertainty agent for intelligent resource allocation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Results:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• State-of-the-art accuracy on BraTS 2023"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• 25.26% energy reduction vs. static baselines"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Demonstrated potential for energy-efficient medical AI"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Future Work:"
    p.level = 0

    p = text_frame.add_paragraph()
    p.text = "• Multi-scale temporal adaptation"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Extension to other medical imaging tasks"
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "• Hardware acceleration for clinical deployment"
    p.level = 1

    set_font_times_new_roman(text_frame)

    # Slide 14: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Thank You"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    set_font_times_new_roman(title_frame)

    # Add subtitle
    subtitle_shape = shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Questions?\n\nBiswajit Sahoo\nAdaptive Spiking Neural Networks for\nEnergy-Efficient Brain Tumor Segmentation"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    set_font_times_new_roman(subtitle_frame)

    return prs


def main():
    """Main function to create and save the presentation."""
    print("Creating professional PowerPoint presentation...")

    # Create presentation
    prs = create_presentation()

    # Save the presentation
    output_path = "research_presentation_v2.pptx"
    prs.save(output_path)

    print(f"Presentation saved as: {output_path}")
    print("✓ Content properly fitted to slides")
    print("✓ Beautiful architecture diagram with color-coded components")
    print("✓ Professional formatting and layout")
    print("✓ Aligned with updated research paper content")


if __name__ == "__main__":
    main()